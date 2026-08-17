#!/usr/bin/env python3
"""Offline smoke and golden-workload checks for the Masterino runtime."""

from __future__ import annotations

import argparse
import importlib
import json
import shutil
import subprocess
import tempfile
import time
from pathlib import Path


IMPORTS = (
    "charset_normalizer",
    "cv2",
    "docx",
    "duckdb",
    "fastexcel",
    "matplotlib",
    "numpy",
    "openpyxl",
    "pandas",
    "PIL",
    "plotly",
    "polars",
    "pptx",
    "pyarrow",
    "pypdf",
    "rapidfuzz",
    "reportlab",
    "scipy",
    "seaborn",
    "sklearn",
    "statsmodels",
    "tabulate",
    "xlsxwriter",
)

TOOLS = ("fc-match", "libreoffice", "officecli", "pandoc", "pdftoppm", "python3")
FONTS = ("Carlito", "Caladea", "Noto Sans CJK SC")


def run(arguments: list[str], timeout: int = 180) -> subprocess.CompletedProcess[str]:
    result = subprocess.run(
        arguments,
        capture_output=True,
        check=False,
        text=True,
        timeout=timeout,
    )
    if result.returncode != 0:
        message = result.stderr.strip() or result.stdout.strip() or "command failed"
        raise RuntimeError(f"{' '.join(arguments)}: {message}")
    return result


def verify_imports() -> dict[str, str]:
    versions: dict[str, str] = {}
    for module_name in IMPORTS:
        module = importlib.import_module(module_name)
        versions[module_name] = str(getattr(module, "__version__", "available"))
    return versions


def verify_tools_and_fonts() -> dict[str, str]:
    tools: dict[str, str] = {}
    for tool in TOOLS:
        path = shutil.which(tool)
        if not path:
            raise RuntimeError(f"required tool is missing: {tool}")
        tools[tool] = path

    for font in FONTS:
        matched = run(["fc-match", "--format", "%{family}", font], timeout=30).stdout
        if font.lower() not in matched.lower():
            raise RuntimeError(f"required font is missing or substituted unexpectedly: {font}")

    return tools


def write_artifacts(directory: Path, rows: int) -> dict[str, int | float | str]:
    import duckdb
    import matplotlib

    matplotlib.use("Agg")

    import matplotlib.pyplot as plt
    import numpy as np
    import pandas as pd
    import polars as pl
    import seaborn as sns
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches
    from pypdf import PdfReader
    from reportlab.pdfgen import canvas

    frame = pd.DataFrame(
        {
            "category": np.resize(np.array(["A", "B", "C", "D"]), rows),
            "value": np.arange(rows, dtype=np.int64),
        }
    )
    expected_sum = int(frame["value"].sum())

    gbk_path = directory / "中文数据.csv"
    frame.head(20).to_csv(gbk_path, index=False, encoding="gbk")
    if len(pd.read_csv(gbk_path, encoding="gbk")) != 20:
        raise RuntimeError("GBK CSV round-trip failed")

    parquet_path = directory / "analysis.parquet"
    frame.to_parquet(parquet_path, index=False)
    polars_sum = int(pl.read_parquet(parquet_path)["value"].sum())
    duckdb_sum = int(duckdb.sql(f"SELECT sum(value) FROM read_parquet('{parquet_path}')").fetchone()[0])
    if polars_sum != expected_sum or duckdb_sum != expected_sum:
        raise RuntimeError("pandas/Polars/DuckDB result mismatch")

    workbook_path = directory / "analysis.xlsx"
    with pd.ExcelWriter(workbook_path, engine="xlsxwriter") as writer:
        frame.head(1000).to_excel(writer, index=False, sheet_name="数据")
    if len(pd.read_excel(workbook_path, engine="openpyxl")) != min(rows, 1000):
        raise RuntimeError("Excel round-trip failed")

    chart_path = directory / "趋势图.png"
    sns.set_theme(style="whitegrid")
    plt.rcParams["font.family"] = "Noto Sans CJK SC"
    sample = frame.groupby("category", as_index=False)["value"].mean()
    figure, axis = plt.subplots(figsize=(8, 4.5))
    sns.barplot(data=sample, x="category", y="value", ax=axis)
    axis.set_title("中文数据趋势")
    figure.tight_layout()
    figure.savefig(chart_path, dpi=144)
    plt.close(figure)

    presentation_path = directory / "python-fallback.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "数据分析报告"
    slide.shapes.add_picture(str(chart_path), Inches(1.5), Inches(1.8), width=Inches(10.3))
    slide.notes_slide.notes_text_frame.text = "Python fallback golden file"
    presentation.save(presentation_path)

    document_path = directory / "report.docx"
    document = Document()
    document.add_heading("数据分析报告", level=1)
    document.add_paragraph(f"记录数：{rows}；合计：{expected_sum}")
    document.save(document_path)

    pdf_path = directory / "report.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(72, 780, f"Rows: {rows}; sum: {expected_sum}")
    pdf.save()
    if len(PdfReader(pdf_path).pages) != 1:
        raise RuntimeError("PDF generation failed")

    return {
        "duckdbSum": duckdb_sum,
        "expectedSum": expected_sum,
        "files": len(list(directory.iterdir())),
        "polarsSum": polars_sum,
        "rows": rows,
    }


def verify_officecli(directory: Path, visual: bool) -> dict[str, str]:
    deck = directory / "officecli-golden.pptx"
    run(["officecli", "create", str(deck)])
    run(
        [
            "officecli",
            "add",
            str(deck),
            "/",
            "--type",
            "slide",
            "--prop",
            "title=Masterino Runtime",
            "--prop",
            "background=FFFFFF",
        ]
    )
    run(["officecli", "validate", str(deck)])
    run(["officecli", "view", str(deck), "outline"])
    run(["officecli", "view", str(deck), "issues", "--json"])

    if visual:
        preview = directory / "officecli-golden.png"
        run(["officecli", "view", str(deck), "screenshot", "--page", "1", "-o", str(preview)])
        if not preview.is_file() or preview.stat().st_size == 0:
            raise RuntimeError("OfficeCLI screenshot was not generated")

    return {"deck": str(deck), "visual": str(visual).lower()}


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--golden", action="store_true")
    parser.add_argument("--output-dir")
    parser.add_argument("--quick", action="store_true")
    parser.add_argument("--rows", type=int, default=1_000_000)
    args = parser.parse_args()

    if args.rows < 1 or args.rows > 5_000_000:
        raise ValueError("rows must be between 1 and 5,000,000")

    started = time.monotonic()
    imports = verify_imports()
    tools = verify_tools_and_fonts()

    requested_rows = args.rows if args.golden else 1_000
    if args.output_dir:
        output_dir = Path(args.output_dir).resolve()
        output_dir.mkdir(parents=True, exist_ok=True)
        cleanup = None
    else:
        cleanup = tempfile.TemporaryDirectory(prefix="masterino-runtime-")
        output_dir = Path(cleanup.name)

    try:
        artifacts = write_artifacts(output_dir, requested_rows)
        office = verify_officecli(output_dir, visual=args.golden and not args.quick)
        result = {
            "artifacts": artifacts,
            "durationMs": round((time.monotonic() - started) * 1000),
            "imports": imports,
            "office": office,
            "outputDirectory": str(output_dir) if args.output_dir else None,
            "success": True,
            "tools": tools,
        }
        print(json.dumps(result, ensure_ascii=False, sort_keys=True))
    finally:
        if cleanup:
            cleanup.cleanup()


if __name__ == "__main__":
    main()
